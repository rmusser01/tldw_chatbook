# Document_Processing_Lib.py
#########################################
# Library to hold functions for ingesting document files
# Supports: DOCX, ODT, RTF, PPTX, XLSX and other office formats
#
####################
# Function List
#
# 1. process_document(file_path, title, author, keywords, chunk_method, chunk_size, chunk_overlap, max_chunk_size)
# 2. process_docx(file_path, title, author, keywords)
# 3. process_with_docling(file_path, title, author, keywords)
# 4. extract_text_from_document(file_path, method='auto')
#
####################
# Import necessary libraries
import os
import re
import gc
from datetime import datetime
from typing import Dict, Any, Optional, List, Union, Tuple
from pathlib import Path
#
# Import External Libs
try:
    from docling.document_converter import DocumentConverter
    from docling.datamodel.base_models import InputFormat
    from docling.datamodel.pipeline_options import PdfPipelineOptions
    from docling.backend.md_backend import MarkdownDocumentBackend
    DOCLING_AVAILABLE = True
except ImportError:
    DOCLING_AVAILABLE = False
    
try:
    import docx
    PYTHON_DOCX_AVAILABLE = True
except ImportError:
    PYTHON_DOCX_AVAILABLE = False

try:
    from odf import text, teletype
    from odf.opendocument import load as load_odf
    ODT_AVAILABLE = True
except ImportError:
    ODT_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

try:
    import openpyxl
    XLSX_AVAILABLE = True
except ImportError:
    XLSX_AVAILABLE = False

try:
    from striprtf.striprtf import rtf_to_text
    RTF_AVAILABLE = True
except ImportError:
    RTF_AVAILABLE = False
#
# Import Local
from ..config import get_cli_setting
from ..LLM_Calls.Summarization_General_Lib import analyze
from ..Metrics.metrics_logger import log_counter, log_histogram
from ..Utils.optional_deps import get_safe_import
from loguru import logger
#
# Constants
# Get media processing config from CLI settings
media_config = {
    'document': {
        'chunk_method': get_cli_setting('media_processing.document', 'chunk_method', 'sentences'),
        'chunk_size': get_cli_setting('media_processing.document', 'chunk_size', 1500),
        'chunk_overlap': get_cli_setting('media_processing.document', 'chunk_overlap', 100),
        'max_summary_length': get_cli_setting('media_processing.document', 'max_summary_length', 1000),
    }
}

# Supported document formats
SUPPORTED_FORMATS = {
    '.docx': 'Microsoft Word',
    '.doc': 'Microsoft Word (Legacy)',
    '.odt': 'OpenDocument Text',
    '.rtf': 'Rich Text Format',
    '.pptx': 'Microsoft PowerPoint',
    '.ppt': 'Microsoft PowerPoint (Legacy)',
    '.xlsx': 'Microsoft Excel',
    '.xls': 'Microsoft Excel (Legacy)',
    '.ods': 'OpenDocument Spreadsheet',
    '.odp': 'OpenDocument Presentation'
}

#######################################################################################################################
# Function Definitions
#

def process_document(
    file_path: str,
    title_override: Optional[str] = None,
    author_override: Optional[str] = None,
    keywords: Optional[List[str]] = None,
    custom_prompt: Optional[str] = None,
    system_prompt: Optional[str] = None,
    summary: Optional[str] = None,
    auto_summarize: bool = False,
    api_name: Optional[str] = None,
    api_key: Optional[str] = None,
    chunk_options: Optional[Dict[str, Any]] = None,
    processing_method: str = 'auto'
) -> Dict[str, Any]:
    """
    Process a document file and extract its content.
    
    Args:
        file_path: Path to the document file
        title_override: Optional title to use instead of extracted title
        author_override: Optional author to use instead of extracted author
        keywords: List of keywords to associate with the document
        custom_prompt: Custom prompt for analysis/summarization
        system_prompt: System prompt for LLM
        summary: Pre-provided summary (if any)
        auto_summarize: Whether to automatically summarize the content
        api_name: LLM API to use for summarization
        api_key: API key for the LLM service
        chunk_options: Chunking configuration
        processing_method: Method to use ('auto', 'docling', 'native')
        
    Returns:
        Dict containing:
            - content: Extracted text content
            - title: Document title
            - author: Document author
            - metadata: Additional metadata
            - summary: Summary (if requested)
            - extraction_successful: Success status
    """
    logger.info(f"Processing document: {file_path}")
    log_counter("document_processing_attempt", labels={"file_path": file_path})
    
    file_ext = Path(file_path).suffix.lower()
    
    if file_ext not in SUPPORTED_FORMATS:
        logger.error(f"Unsupported document format: {file_ext}")
        return {
            'content': '',
            'title': os.path.basename(file_path),
            'author': 'Unknown',
            'metadata': {'error': f'Unsupported format: {file_ext}'},
            'extraction_successful': False
        }
    
    # Default chunk options
    if chunk_options is None:
        chunk_options = {
            'method': media_config['document']['chunk_method'],
            'max_size': media_config['document']['chunk_size'],
            'overlap': media_config['document']['chunk_overlap'],
            'adaptive': True,
            'language': 'en'
        }
    
    try:
        # Determine processing method
        if processing_method == 'auto':
            if DOCLING_AVAILABLE:
                processing_method = 'docling'
            else:
                processing_method = 'native'
        
        # Extract content based on method
        if processing_method == 'docling' and DOCLING_AVAILABLE:
            result = process_with_docling(file_path, title_override, author_override, keywords)
        else:
            # Use native libraries based on file type
            if file_ext == '.docx' and PYTHON_DOCX_AVAILABLE:
                result = process_docx(file_path, title_override, author_override, keywords)
            elif file_ext == '.odt' and ODT_AVAILABLE:
                result = process_odt(file_path, title_override, author_override, keywords)
            elif file_ext == '.rtf' and RTF_AVAILABLE:
                result = process_rtf(file_path, title_override, author_override, keywords)
            elif file_ext == '.pptx' and PPTX_AVAILABLE:
                result = process_pptx(file_path, title_override, author_override, keywords)
            elif file_ext in ['.xlsx', '.xls'] and XLSX_AVAILABLE:
                result = process_xlsx(file_path, title_override, author_override, keywords)
            else:
                # Fallback to Docling if available
                if DOCLING_AVAILABLE:
                    result = process_with_docling(file_path, title_override, author_override, keywords)
                else:
                    return {
                        'content': '',
                        'title': os.path.basename(file_path),
                        'author': 'Unknown',
                        'metadata': {'error': f'No parser available for {file_ext}'},
                        'extraction_successful': False
                    }
        
        # Add summarization if requested
        if auto_summarize and api_name and result.get('extraction_successful'):
            content = result.get('content', '')
            if content:
                if custom_prompt:
                    summary_prompt = custom_prompt
                else:
                    summary_prompt = f"Please provide a comprehensive summary of this {SUPPORTED_FORMATS.get(file_ext, 'document')}."
                
                summary = analyze(
                    input_data=content,
                    custom_prompt_arg=summary_prompt,
                    api_name=api_name,
                    api_key=api_key,
                    temp=0.7,
                    system_message=system_prompt
                )
                result['summary'] = summary
            else:
                result['summary'] = "No content available to summarize."
        elif summary:
            result['summary'] = summary
        
        log_counter("document_processing_success", labels={"method": processing_method, "format": file_ext})
        return result
        
    except Exception as e:
        logger.error(f"Error processing document {file_path}: {str(e)}", exc_info=True)
        log_counter("document_processing_error", labels={"error": str(e), "format": file_ext})
        return {
            'content': '',
            'title': os.path.basename(file_path),
            'author': 'Unknown',
            'metadata': {'error': str(e)},
            'extraction_successful': False
        }


def process_with_docling(
    file_path: str,
    title_override: Optional[str] = None,
    author_override: Optional[str] = None,
    keywords: Optional[List[str]] = None
) -> Dict[str, Any]:
    """
    Process a document using Docling.
    
    Docling provides advanced document understanding and can handle
    complex layouts, tables, and various document formats.
    """
    logger.info(f"Processing document with Docling: {file_path}")
    
    try:
        # Initialize Docling converter
        converter = DocumentConverter()
        
        # Configure pipeline options
        pipeline_options = PdfPipelineOptions()
        pipeline_options.do_ocr = False  # Disable OCR for documents (not PDFs)
        pipeline_options.do_table_structure = True  # Extract tables
        
        # Convert document
        result = converter.convert(file_path, pipeline_options=pipeline_options)
        
        # Export to markdown
        md_backend = MarkdownDocumentBackend()
        markdown_content = result.document.export_to_markdown(backend=md_backend)
        
        # Extract metadata
        metadata = {
            'num_pages': len(result.document.pages) if hasattr(result.document, 'pages') else 1,
            'processing_method': 'docling',
            'tables_found': len(result.document.tables) if hasattr(result.document, 'tables') else 0,
        }
        
        # Extract title and author from document if available
        title = title_override or result.document.title if hasattr(result.document, 'title') else os.path.basename(file_path)
        author = author_override or result.document.author if hasattr(result.document, 'author') else 'Unknown'
        
        return {
            'content': markdown_content,
            'title': title,
            'author': author,
            'metadata': metadata,
            'keywords': keywords or [],
            'extraction_successful': True
        }
        
    except Exception as e:
        logger.error(f"Error processing with Docling: {str(e)}", exc_info=True)
        raise


def process_docx(
    file_path: str,
    title_override: Optional[str] = None,
    author_override: Optional[str] = None,
    keywords: Optional[List[str]] = None
) -> Dict[str, Any]:
    """Process a DOCX file using python-docx."""
    logger.info(f"Processing DOCX file: {file_path}")
    
    try:
        doc = docx.Document(file_path)
        
        # Extract text content
        content_parts = []
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                content_parts.append(paragraph.text)
        
        # Extract tables
        for table in doc.tables:
            table_text = []
            for row in table.rows:
                row_text = []
                for cell in row.cells:
                    cell_text = cell.text.strip()
                    if cell_text:
                        row_text.append(cell_text)
                if row_text:
                    table_text.append(' | '.join(row_text))
            if table_text:
                content_parts.append('\n' + '\n'.join(table_text) + '\n')
        
        content = '\n\n'.join(content_parts)
        
        # Extract metadata
        core_props = doc.core_properties
        title = title_override or core_props.title or os.path.basename(file_path)
        author = author_override or core_props.author or 'Unknown'
        
        metadata = {
            'created': str(core_props.created) if core_props.created else None,
            'modified': str(core_props.modified) if core_props.modified else None,
            'subject': core_props.subject,
            'processing_method': 'python-docx'
        }
        
        return {
            'content': content,
            'title': title,
            'author': author,
            'metadata': metadata,
            'keywords': keywords or [],
            'extraction_successful': True
        }
        
    except Exception as e:
        logger.error(f"Error processing DOCX: {str(e)}", exc_info=True)
        raise


def process_odt(
    file_path: str,
    title_override: Optional[str] = None,
    author_override: Optional[str] = None,
    keywords: Optional[List[str]] = None
) -> Dict[str, Any]:
    """Process an ODT file using odfpy."""
    logger.info(f"Processing ODT file: {file_path}")
    
    try:
        doc = load_odf(file_path)
        
        # Extract text content
        all_paragraphs = doc.getElementsByType(text.P)
        content_parts = []
        
        for paragraph in all_paragraphs:
            text_content = teletype.extractText(paragraph).strip()
            if text_content:
                content_parts.append(text_content)
        
        content = '\n\n'.join(content_parts)
        
        # Extract metadata
        meta = doc.meta
        title = title_override or (meta.getElementsByType(text.Title)[0].firstChild.data if meta.getElementsByType(text.Title) else os.path.basename(file_path))
        author = author_override or 'Unknown'  # ODT metadata extraction is limited
        
        return {
            'content': content,
            'title': title,
            'author': author,
            'metadata': {'processing_method': 'odfpy'},
            'keywords': keywords or [],
            'extraction_successful': True
        }
        
    except Exception as e:
        logger.error(f"Error processing ODT: {str(e)}", exc_info=True)
        raise


def process_rtf(
    file_path: str,
    title_override: Optional[str] = None,
    author_override: Optional[str] = None,
    keywords: Optional[List[str]] = None
) -> Dict[str, Any]:
    """Process an RTF file using striprtf."""
    logger.info(f"Processing RTF file: {file_path}")
    
    try:
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as file:
            rtf_content = file.read()
        
        # Convert RTF to plain text
        text_content = rtf_to_text(rtf_content)
        
        return {
            'content': text_content,
            'title': title_override or os.path.basename(file_path),
            'author': author_override or 'Unknown',
            'metadata': {'processing_method': 'striprtf'},
            'keywords': keywords or [],
            'extraction_successful': True
        }
        
    except Exception as e:
        logger.error(f"Error processing RTF: {str(e)}", exc_info=True)
        raise


def process_pptx(
    file_path: str,
    title_override: Optional[str] = None,
    author_override: Optional[str] = None,
    keywords: Optional[List[str]] = None
) -> Dict[str, Any]:
    """Process a PPTX file using python-pptx."""
    logger.info(f"Processing PPTX file: {file_path}")
    
    try:
        prs = Presentation(file_path)
        
        content_parts = []
        
        # Extract text from slides
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = [f"## Slide {slide_num}"]
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text)
                
                # Extract text from tables
                if shape.has_table:
                    table_text = []
                    for row in shape.table.rows:
                        row_text = []
                        for cell in row.cells:
                            if cell.text.strip():
                                row_text.append(cell.text.strip())
                        if row_text:
                            table_text.append(' | '.join(row_text))
                    if table_text:
                        slide_text.append('\n'.join(table_text))
            
            if len(slide_text) > 1:  # More than just the slide header
                content_parts.append('\n'.join(slide_text))
        
        content = '\n\n'.join(content_parts)
        
        # Extract metadata
        core_props = prs.core_properties
        title = title_override or core_props.title or os.path.basename(file_path)
        author = author_override or core_props.author or 'Unknown'
        
        metadata = {
            'num_slides': len(prs.slides),
            'created': str(core_props.created) if core_props.created else None,
            'modified': str(core_props.modified) if core_props.modified else None,
            'processing_method': 'python-pptx'
        }
        
        return {
            'content': content,
            'title': title,
            'author': author,
            'metadata': metadata,
            'keywords': keywords or [],
            'extraction_successful': True
        }
        
    except Exception as e:
        logger.error(f"Error processing PPTX: {str(e)}", exc_info=True)
        raise


def process_xlsx(
    file_path: str,
    title_override: Optional[str] = None,
    author_override: Optional[str] = None,
    keywords: Optional[List[str]] = None
) -> Dict[str, Any]:
    """Process an XLSX file using openpyxl."""
    logger.info(f"Processing XLSX file: {file_path}")
    
    try:
        wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
        
        content_parts = []
        
        # Process each sheet
        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            sheet_content = [f"## Sheet: {sheet_name}"]
            
            # Extract data from cells
            rows_data = []
            for row in sheet.iter_rows(values_only=True):
                # Filter out empty rows
                row_values = [str(cell) if cell is not None else '' for cell in row]
                if any(val.strip() for val in row_values):
                    rows_data.append(' | '.join(row_values))
            
            if rows_data:
                sheet_content.extend(rows_data[:100])  # Limit to first 100 rows
                if len(rows_data) > 100:
                    sheet_content.append(f"... and {len(rows_data) - 100} more rows")
                
                content_parts.append('\n'.join(sheet_content))
        
        content = '\n\n'.join(content_parts)
        
        # Extract metadata
        title = title_override or wb.properties.title or os.path.basename(file_path)
        author = author_override or wb.properties.creator or 'Unknown'
        
        metadata = {
            'num_sheets': len(wb.sheetnames),
            'created': str(wb.properties.created) if wb.properties.created else None,
            'modified': str(wb.properties.modified) if wb.properties.modified else None,
            'processing_method': 'openpyxl'
        }
        
        wb.close()
        
        return {
            'content': content,
            'title': title,
            'author': author,
            'metadata': metadata,
            'keywords': keywords or [],
            'extraction_successful': True
        }
        
    except Exception as e:
        logger.error(f"Error processing XLSX: {str(e)}", exc_info=True)
        raise


def extract_text_from_document(file_path: str, method: str = 'auto') -> str:
    """
    Simple text extraction from a document file.
    
    Args:
        file_path: Path to the document
        method: Extraction method ('auto', 'docling', 'native')
        
    Returns:
        Extracted text as string
    """
    result = process_document(file_path, processing_method=method)
    
    if result.get('extraction_successful'):
        return result.get('content', '')
    else:
        raise ValueError(f"Failed to extract text: {result.get('metadata', {}).get('error', 'Unknown error')}")


# Compatibility function for existing code
def process_document_with_docling(file_path: str, **kwargs) -> Dict[str, Any]:
    """Compatibility wrapper that forces Docling processing."""
    return process_document(file_path, processing_method='docling', **kwargs)